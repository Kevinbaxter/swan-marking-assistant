import os
import re
from collections import Counter

import streamlit as st
from docx import Document
from pptx import Presentation
import pandas as pd

# -----------------------------
# Config
# -----------------------------
SHORT_PARAGRAPH_LIMIT = 260

st.set_page_config(
    page_title="🦢 SWAN Marking Assistant",
    page_icon="🦢",
    layout="centered"
)

# -----------------------------
# Header
# -----------------------------
st.markdown(
    """
    <div style='text-align:center; margin-bottom: 20px;'>
        <h1 style='color:#ffffff;'>SWAN Marking Assistant</h1>
        <p style='font-size:18px; color:#cccccc;'>
            Upload your writing and get Strengths, Weaknesses, Actions and Next Steps,
            written in a friendly, teacher-style voice.
        </p>
    </div>
    """,
    unsafe_allow_html=True
)

# -----------------------------
# Extraction helpers
# -----------------------------
def extract_text_from_docx(file):
    doc = Document(file)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return doc, paragraphs

def extract_text_from_xlsx(file):
    xls = pd.ExcelFile(file)
    text_blocks = []
    for sheet in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet)
        text_blocks.append(f"Sheet: {sheet}")
        flat_values = df.astype(str).fillna("").values.flatten().tolist()
        text_blocks.extend([v for v in flat_values if v.strip()])
    return None, text_blocks

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text_blocks = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            text_blocks.append(f"Slide {slide_idx}:")
            text_blocks.extend(slide_text)
    return None, text_blocks

def extract_text(file, ext):
    if ext == ".docx":
        return extract_text_from_docx(file)
    elif ext == ".xlsx":
        return extract_text_from_xlsx(file)
    elif ext == ".pptx":
        return extract_text_from_pptx(file)
    else:
        return None, []

# -----------------------------
# Low-level analysis helpers
# -----------------------------
def count_headings_docx(doc):
    if doc is None:
        return 0
    return sum(1 for p in doc.paragraphs if p.style and p.style.name.startswith("Heading"))

def has_bullets_docx(doc):
    if doc is None:
        return False
    return any("List" in (p.style.name if p.style else "") for p in doc.paragraphs)

def find_short_paragraphs(paragraphs):
    return [p for p in paragraphs if len(p) < SHORT_PARAGRAPH_LIMIT]

def sentence_lengths(text):
    sentences = re.split(r"[.!?]", text)
    lengths = [len(s.split()) for s in sentences if len(s.split()) > 0]
    return lengths

def vocab_stats(text):
    words = re.findall(r"\b[a-zA-Z]{3,}\b", text.lower())
    if not words:
        return 0.0, 0
    unique = len(set(words))
    total = len(words)
    return unique / total, total

def detect_tone(text):
    t = text.lower()
    informal_markers = ["gonna", "wanna", "yeah", "kinda", "sort of", "like,"]
    formal_markers = ["therefore", "however", "in conclusion", "moreover", "furthermore"]
    informal = any(w in t for w in informal_markers)
    formal = any(w in t for w in formal_markers)
    if formal and not informal:
        return "mostly formal"
    if informal and not formal:
        return "quite informal"
    if formal and informal:
        return "a mix of formal and informal"
    return "neutral"

def estimate_cefr(avg_len, vocab_ratio):
    # Very rough, gentle estimate
    if avg_len < 8 and vocab_ratio < 0.25:
        return "around A2 level"
    if 8 <= avg_len <= 18 and 0.25 <= vocab_ratio <= 0.4:
        return "around B1 level"
    if avg_len > 18 and vocab_ratio > 0.35:
        return "moving towards B2 level"
    return "somewhere between A2 and B1 level"

# -----------------------------
# Main SWAN student analysis
# -----------------------------
def analyse_student_writing(doc, paragraphs, ext):
    strengths = []
    weaknesses = []
    actions = []
    next_steps = []

    if not paragraphs:
        weaknesses.append("There wasn’t any clear text to read in your file, so it’s hard to comment on your writing.")
        next_steps.append("Try uploading a version that includes your full writing, not just a blank template or image.")
        summary = "I couldn’t really see your writing this time, so I can’t give proper feedback yet."
        metrics = {}
        return strengths, weaknesses, actions, next_steps, summary, metrics

    text = " ".join(paragraphs)
    lower_text = text.lower()

    # STRUCTURE (Word only)
    if ext == ".docx":
        if count_headings_docx(doc) >= 1:
            strengths.append("You’ve started to organise your work with headings, which helps the reader follow your ideas.")
        else:
            weaknesses.append("At the moment, your work doesn’t really use headings to guide the reader.")
            actions.append("Try adding simple headings to show where each new idea or section begins.")

        if has_bullets_docx(doc):
            strengths.append("You use bullet points in places, which can make key information stand out clearly.")
        else:
            weaknesses.append("Some of your ideas are in long blocks of text and could be clearer as bullet points.")
            actions.append("Where you list examples or points, try using bullet points to make them easier to read.")

    # PARAGRAPH DEVELOPMENT
    short_paras = find_short_paragraphs(paragraphs)
    if not short_paras:
        strengths.append("Your paragraphs generally have enough detail to explain your ideas clearly.")
    else:
        weaknesses.append("Some of your paragraphs are quite short and feel like they stop before the idea is fully explained.")
        actions.append("Choose one short paragraph and add an extra sentence that gives an example or explains your point more.")

    # CONCLUSION
    last_para = paragraphs[-1].strip().lower()
    if any(p in last_para for p in ["in conclusion", "overall", "to sum up", "in summary"]):
        strengths.append("You’ve tried to round off your writing with a concluding idea, which helps give it a clear ending.")
    else:
        weaknesses.append("Your writing finishes quite suddenly without a clear final sentence to bring your ideas together.")
        actions.append("Add a short final sentence that sums up your main point or how you feel about the topic.")

    # SENTENCE VARIETY
    lengths = sentence_lengths(text)
    if lengths:
        avg_len = sum(lengths) / len(lengths)
        if avg_len < 9:
            weaknesses.append("A lot of your sentences are very short, which can make the writing feel a bit choppy.")
            actions.append("Try joining two short sentences together using a linking word like 'because', 'so' or 'which'.")
        elif avg_len > 24:
            weaknesses.append("Some of your sentences are quite long, which can make them harder to follow.")
            actions.append("Pick one long sentence and see if you can split it into two shorter ones without losing meaning.")
        else:
            strengths.append("You use a mix of shorter and longer sentences, which helps your writing flow more naturally.")
    else:
        avg_len = 0

    # VOCABULARY
    vocab_ratio, total_words = vocab_stats(text)
    if total_words > 0:
        if vocab_ratio > 0.4:
            strengths.append("You’re beginning to use a good range of vocabulary, which makes your writing more interesting to read.")
        elif vocab_ratio < 0.25:
            weaknesses.append("You repeat some words quite a lot, which can make the writing feel a bit limited.")
            actions.append("Choose one common word you use a lot and try replacing it with a different word in one or two places.")

    # LINKING WORDS
    LINKERS = [
        "however", "therefore", "in addition", "furthermore", "moreover",
        "for example", "for instance", "as a result", "on the other hand"
    ]
    if any(l in lower_text for l in LINKERS):
        strengths.append("You’ve started to use linking phrases to connect your ideas, which helps the reader follow your thinking.")
    else:
        weaknesses.append("Your ideas sometimes feel like separate points rather than a joined-up piece of writing.")
        actions.append("Try adding phrases like 'for example', 'as a result' or 'in addition' to show how your ideas connect.")

    # ARGUMENT / EXPLANATION
    ARG_MARKERS = ["because", "this shows", "this suggests", "so that", "therefore", "as a result"]
    if any(m in lower_text for m in ARG_MARKERS):
        strengths.append("You don’t just make points – you also try to explain or justify them, which is a really positive skill.")
    else:
        weaknesses.append("Sometimes you make a point but don’t fully explain why it matters or what it shows.")
        actions.append("After making a point, add a short phrase like 'this shows that…' or 'this is important because…'.")

    # METRICS / GENTLE SCORES
    tone = detect_tone(text)
    cefr = estimate_cefr(avg_len, vocab_ratio)
    # Soft, gentle 1–5 style feelings
    clarity_score = 3
    variety_score = 3
    vocab_score = 3

    if avg_len and 10 <= avg_len <= 20:
        clarity_score = 4
    if avg_len and (avg_len < 7 or avg_len > 26):
        clarity_score = 2

    if lengths and max(lengths) - min(lengths) > 10:
        variety_score = 4
    elif lengths and max(lengths) - min(lengths) < 5:
        variety_score = 2

    if vocab_ratio > 0.4:
        vocab_score = 4
    elif vocab_ratio < 0.25:
        vocab_score = 2

    metrics = {
        "clarity": clarity_score,
        "variety": variety_score,
        "vocab": vocab_score,
        "tone": tone,
        "cefr": cefr
    }

    # OVERALL SUMMARY (Style L – short paragraph)
    summary_parts = []
    summary_parts.append(
        "You’ve made a thoughtful start here, and it’s clear you’re trying to express your ideas in a clear and organised way."
    )
    if strengths:
        summary_parts.append(
            "Some parts of your writing already work well, especially where you explain your points or use structure to guide the reader."
        )
    if weaknesses:
        summary_parts.append(
            "At times, your ideas feel a little brief or jump from one point to another, but this is something that improves naturally with practice."
        )
    if vocab_score >= 3:
        summary_parts.append(
            "Your vocabulary choices show potential, and you’re beginning to experiment with different ways of expressing yourself."
        )
    summary_parts.append(
        "With a bit more detail and a focus on linking your ideas smoothly, your writing will continue to grow in confidence."
    )
    summary = " ".join(summary_parts)

    # NEXT STEPS – always a few gentle, practical ideas
    next_steps.append("Read your work aloud slowly and check that each sentence flows naturally into the next one.")
    next_steps.append("Choose one paragraph and see if you can add one extra sentence that gives an example or explains your point more clearly.")
    next_steps.append("Look for one place where you can add a linking phrase like 'for example', 'as a result' or 'in addition' to join ideas together.")

    return strengths, weaknesses, actions, next_steps, summary, metrics

# -----------------------------
# UI
# -----------------------------
uploaded = st.file_uploader(
    "Upload a document (.docx, .xlsx, .pptx)",
    type=["docx", "xlsx", "pptx"]
)

if st.button("🔄 Reset"):
    st.session_state.clear()
    st.rerun()

if uploaded:
    ext = os.path.splitext(uploaded.name)[1].lower()
    st.info(f"File detected: **{uploaded.name}** ({ext})")

    doc, paragraphs = extract_text(uploaded, ext)
    strengths, weaknesses, actions, next_steps, summary, metrics = analyse_student_writing(doc, paragraphs, ext)

    st.subheader("Overall summary")
    st.write(summary)

    st.subheader("Strengths")
    if strengths:
        for s in strengths[:8]:
            st.write("•", s)
    else:
        st.write("I couldn’t see clear strengths this time, mainly because there wasn’t enough continuous writing to analyse.")

    st.subheader("Things to keep working on")
    if weaknesses:
        for w in weaknesses[:8]:
            st.write("•", w)
    else:
        st.write("There aren’t any major issues that stand out strongly from this basic check.")

    st.subheader("Actions you can take")
    if actions:
        for a in actions[:8]:
            st.write("•", a)
    else:
        st.write("There aren’t any specific action points beyond the general next steps below.")

    st.subheader("Next steps")
    for n in next_steps[:8]:
        st.write("•", n)

    st.subheader("Gentle deeper insights")
    if metrics:
        st.write(f"- Your clarity feels like it’s around about **{metrics['clarity']} out of 5** at the moment.")
        st.write(f"- Your sentence variety feels around **{metrics['variety']} out of 5**, with room to experiment a bit more.")
        st.write(f"- Your vocabulary range feels around **{metrics['vocab']} out of 5**, with good potential to grow.")
        st.write(f"- The tone of your writing comes across as **{metrics['tone']}**.")
        st.write(f"- Overall, your writing feels **{metrics['cefr']}**, based on sentence length and vocabulary.")
    else:
        st.write("There wasn’t enough clear text to give deeper insights this time.")

    st.divider()

    report_text = (
        f"SWAN Feedback Report: {uploaded.name}\n"
        + "="*30 + "\n\n"
        + "OVERALL SUMMARY:\n"
        + summary + "\n\n"
        + "STRENGTHS:\n" + "\n".join(f"- {s}" for s in strengths) + "\n\n"
        + "THINGS TO KEEP WORKING ON:\n" + "\n".join(f"- {w}" for w in weaknesses) + "\n\n"
        + "ACTIONS YOU CAN TAKE:\n" + "\n".join(f"- {a}" for a in actions) + "\n\n"
        + "NEXT STEPS:\n" + "\n".join(f"- {n}" for n in next_steps) + "\n\n"
    )

    if metrics:
        report_text += (
            "GENTLE DEEPER INSIGHTS:\n"
            f"- Clarity feels around {metrics['clarity']} out of 5.\n"
            f"- Sentence variety feels around {metrics['variety']} out of 5.\n"
            f"- Vocabulary range feels around {metrics['vocab']} out of 5.\n"
            f"- Tone: {metrics['tone']}.\n"
            f"- Overall level: {metrics['cefr']}.\n"
        )

    st.download_button(
        label="📥 Download feedback as text",
        data=report_text,
        file_name=f"SWAN_Feedback_{uploaded.name}.txt",
        mime="text/plain"
    )

else:
    st.caption("Upload a piece of student writing to begin the SWAN analysis.")