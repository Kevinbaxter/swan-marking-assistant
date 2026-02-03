import os
import re
from collections import Counter

import streamlit as st
from docx import Document
from pptx import Presentation
import pandas as pd

# -----------------------------
# Config
# -----------------------------
SHORT_PARAGRAPH_LIMIT = 260

st.set_page_config(
    page_title="🦢 SWAN Marking Assistant",
    page_icon="🦢",
    layout="centered"
)

# -----------------------------
# Text extraction helpers
# -----------------------------
def extract_text_from_docx(file):
    """Extract paragraphs from a Word document."""
    doc = Document(file)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return doc, paragraphs

def extract_text_from_xlsx(file):
    """Extract text-like content from all sheets in an Excel workbook."""
    xls = pd.ExcelFile(file)
    text_blocks = []

    for sheet in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet)
        text_blocks.append(f"Sheet: {sheet}")
        # Flatten all cell values to strings
        flat_values = df.astype(str).fillna("").values.flatten().tolist()
        text_blocks.extend([v for v in flat_values if v.strip()])

    return None, text_blocks  # No 'doc' object here, just text

def extract_text_from_pptx(file):
    """Extract text from all slides in a PowerPoint presentation."""
    prs = Presentation(file)
    text_blocks = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            text_blocks.append(f"Slide {slide_idx}:")
            text_blocks.extend(slide_text)

    return None, text_blocks  # No 'doc' object here, just text

# -----------------------------
# Structural analysis helpers
# -----------------------------
def count_headings_docx(doc):
    if doc is None:
        return 0
    return sum(1 for p in doc.paragraphs if p.style and p.style.name.startswith("Heading"))

def has_bullets_docx(doc):
    if doc is None:
        return False
    return any("List" in (p.style.name if p.style else "") for p in doc.paragraphs)

def find_short_paragraphs(paragraphs):
    return [p for p in paragraphs if len(p) < SHORT_PARAGRAPH_LIMIT]

def find_spelling_issues(paragraphs):
    text = " ".join(paragraphs).lower()
    words = re.findall(r"\b[a-zA-Z]{3,}\b", text)
    counts = Counter(words)
    # crude heuristic: words repeated 3+ times might be misused/misspelt
    return [w for w, c in counts.items() if c >= 3]

# -----------------------------
# SWAN engine
# -----------------------------
def run_swan_analysis(file, ext):
    """
    Main SWAN analysis dispatcher.
    Returns: strengths, weaknesses, actions, next_steps
    """
    doc = None
    paragraphs = []

    # 1. Extract text based on file type
    if ext == ".docx":
        doc, paragraphs = extract_text_from_docx(file)
    elif ext == ".xlsx":
        doc, paragraphs = extract_text_from_xlsx(file)
    elif ext == ".pptx":
        doc, paragraphs = extract_text_from_pptx(file)
    else:
        return [], ["Unsupported file type."], [], []

    if not paragraphs:
        return [], ["No readable content was found in the file."], [], []

    strengths = []
    weaknesses = []
    actions = []
    next_steps = []

    # -------------------------
    # Structural strengths
    # -------------------------
    if ext == ".docx":
        if count_headings_docx(doc) >= 1:
            strengths.append("You have used headings to organise your work.")
        if has_bullets_docx(doc):
            strengths.append("Bullet points help make your ideas clear and easy to read.")

    # For Excel and PowerPoint, we can infer structure from markers
    if ext in [".xlsx", ".pptx"]:
        if any("sheet:" in p.lower() for p in paragraphs):
            strengths.append("Your work is organised into clear sections or sheets.")
        if any("slide" in p.lower() for p in paragraphs):
            strengths.append("Your slides are clearly separated and labelled.")

    # Paragraph/detail strength
    short_paras = find_short_paragraphs(paragraphs)
    if not short_paras:
        strengths.append("Most of your sections are developed with enough detail.")

    # -------------------------
    # Weaknesses & actions
    # -------------------------
    if short_paras:
        weaknesses.append("Some sections are very short and lack detail.")
        actions.append("Action: Choose one short section and add at least one example or explanation.")

    # Conclusion check (last block)
    if len(paragraphs[-1]) < SHORT_PARAGRAPH_LIMIT:
        weaknesses.append("There is no clear conclusion at the end of your work.")
        actions.append("Action: Add a short conclusion that summarises your key points and links back to the question or purpose.")

    # Spelling / repetition
    spelling_issues = find_spelling_issues(paragraphs)
    if spelling_issues:
        weaknesses.append("The same word or phrase is repeated many times, which may be a spelling or vocabulary issue.")
        actions.append("Action: Review repeated words and check their spelling or replace some with alternatives.")

    # -------------------------
    # Next steps (N in SWAN)
    # -------------------------
    next_steps.append("Read your work aloud to check that it flows logically and makes sense.")
    next_steps.append("Compare your structure to a model answer or exemplar to see how you could improve organisation.")
    next_steps.append("Ask a peer or teacher to highlight one section that could be clearer, then rewrite it.")

    return strengths, weaknesses, actions, next_steps

# -----------------------------
# UI
# -----------------------------
st.title("🦢 SWAN Marking Assistant")

st.markdown(
    "Upload a **Word, Excel, or PowerPoint** file and get structured "
    "**Strengths, Weaknesses, Actions, and Next Steps** feedback."
)

uploaded = st.file_uploader(
    "Upload a document (.docx, .xlsx, .pptx)",
    type=["docx", "xlsx", "pptx"]
)

if uploaded:
    ext = os.path.splitext(uploaded.name)[1].lower()

    st.info(f"File detected: **{uploaded.name}** ({ext})")

    strengths, weaknesses, actions, next_steps = run_swan_analysis(uploaded, ext)

    st.subheader("Strengths")
    if strengths:
        for s in strengths[:5]:
            st.write("•", s)
    else:
        st.write("No clear strengths were detected. Try adding more structure or detail.")

    st.subheader("Weaknesses")
    if weaknesses:
        for w in weaknesses[:5]:
            st.write("•", w)
    else:
        st.write("No major weaknesses were detected from this basic analysis.")

    st.subheader("Actions")
    if actions:
        for a in actions[:5]:
            st.write("•", a)
    else:
        st.write("No specific actions generated.")

    st.subheader("Next Steps")
    if next_steps:
        for n in next_steps[:5]:
            st.write("•", n)
    else:
        st.write("No next steps generated.")
else:
    st.caption("Waiting for a file upload to begin SWAN analysis…")