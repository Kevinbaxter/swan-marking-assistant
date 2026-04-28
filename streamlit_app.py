import os
import re
from collections import Counter

import streamlit as st
from docx import Document
from pptx import Presentation
import pandas as pd

# -----------------------------
# Config
# -----------------------------
SHORT_PARAGRAPH_LIMIT = 260

st.set_page_config(
    page_title="🦢 SWAN Marking Assistant",
    page_icon="🦢",
    layout="centered"
)

# -----------------------------
# Header + Logo
# -----------------------------
st.markdown(
    """
    <div style='text-align:center; margin-bottom: 20px;'>
        <img src='https://raw.githubusercontent.com/Kevinbaxter/swan-marking-assistant/main/kclogo.png' 
             width='300' style='margin-bottom:10px;'/>
        <h1 style='color:#ffffff;'>SWAN Marking Assistant</h1>
        <p style='font-size:18px; color:#cccccc;'>
            Upload a document and receive structured Strengths, Weaknesses, Actions and Next Steps.
        </p>
    </div>
    """,
    unsafe_allow_html=True
)

# -----------------------------
# Extraction helpers
# -----------------------------
def extract_text_from_docx(file):
    doc = Document(file)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return doc, paragraphs

def extract_text_from_xlsx(file):
    xls = pd.ExcelFile(file)
    text_blocks = []

    for sheet in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet)
        text_blocks.append(f"Sheet: {sheet}")
        flat_values = df.astype(str).fillna("").values.flatten().tolist()
        text_blocks.extend([v for v in flat_values if v.strip()])

    return None, text_blocks

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text_blocks = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            text_blocks.append(f"Slide {slide_idx}:")
            text_blocks.extend(slide_text)

    return None, text_blocks

# -----------------------------
# Analysis helpers
# -----------------------------
def count_headings_docx(doc):
    if doc is None:
        return 0
    return sum(1 for p in doc.paragraphs if p.style and p.style.name.startswith("Heading"))

def has_bullets_docx(doc):
    if doc is None:
        return False
    return any("List" in (p.style.name if p.style else "") for p in doc.paragraphs)

def find_short_paragraphs(paragraphs):
    return [p for p in paragraphs if len(p) < SHORT_PARAGRAPH_LIMIT]

def find_spelling_issues(paragraphs):
    text = " ".join(paragraphs).lower()
    words = re.findall(r"\b[a-zA-Z]{3,}\b", text)
    counts = Counter(words)
    return [w for w, c in counts.items() if c >= 3]

# -----------------------------
# SWAN analysis engine
# -----------------------------
def run_swan_analysis(file, ext):
    doc = None
    paragraphs = []

    # Extract text
    if ext == ".docx":
        doc, paragraphs = extract_text_from_docx(file)
    elif ext == ".xlsx":
        doc, paragraphs = extract_text_from_xlsx(file)
    elif ext == ".pptx":
        doc, paragraphs = extract_text_from_pptx(file)
    else:
        return [], ["Unsupported file type."], [], []

    if not paragraphs:
        return [], ["No readable content was found in the file."], [], []

    strengths = []
    weaknesses = []
    actions = []
    next_steps = []

    text = " ".join(paragraphs).lower()

    # STRUCTURE (Word only)
    if ext == ".docx":
        if count_headings_docx(doc) >= 1:
            strengths.append("You have used headings to organise your work clearly.")
        else:
            weaknesses.append("Your work would benefit from clear headings to guide the reader.")
            actions.append("Add headings to show where each new idea or section begins.")

        if has_bullets_docx(doc):
            strengths.append("Bullet points help make your ideas clear and easy to read.")
        else:
            weaknesses.append("Some sections could be clearer with bullet points.")
            actions.append("Use bullet points for lists or key ideas.")

    # PARAGRAPH DEVELOPMENT
    short_paras = find_short_paragraphs(paragraphs)
    if not short_paras:
        strengths.append("Your paragraphs are well-developed with enough detail.")
    else:
        weaknesses.append("Some paragraphs are very short and lack development.")
        actions.append("Choose one short paragraph and expand it with an example or explanation.")

    # CONCLUSION CHECK
    last_para = paragraphs[-1].lower()
    if any(phrase in last_para for phrase in ["in conclusion", "overall", "to sum up", "in summary"]):
        strengths.append("You have included a clear concluding section.")
    else:
        weaknesses.append("Your work ends abruptly without a clear conclusion.")
        actions.append("Add a short conclusion that summarises your key points.")

    # SPELLING / REPETITION
    spelling_issues = find_spelling_issues(paragraphs)
    if spelling_issues:
        weaknesses.append("Some words appear repeatedly and may be misspelt.")
        actions.append("Review repeated words and check their spelling or replace them with alternatives.")

    # SENTENCE VARIETY
    sentences = re.split(r"[.!?]", text)
    sentence_lengths = [len(s.split()) for s in sentences if len(s.split()) > 0]

    if sentence_lengths:
        avg_len = sum(sentence_lengths) / len(sentence_lengths)

        if avg_len < 10:
            weaknesses.append("Many sentences are very short, which makes the writing feel choppy.")
            actions.append("Combine some short sentences to create more complex ones.")
        elif avg_len > 25:
            weaknesses.append("Some sentences are very long and may be hard to follow.")
            actions.append("Split long sentences into two shorter ones to improve clarity.")
        else:
            strengths.append("You use a good mix of short and longer sentences.")

    # VOCABULARY RICHNESS
    words = re.findall(r"\b[a-zA-Z]{3,}\b", text)
    if words:
        unique_ratio = len(set(words)) / len(words)

        if unique_ratio > 0.4:
            strengths.append("Your vocabulary is varied and precise.")
        elif unique_ratio < 0.25:
            weaknesses.append("Your vocabulary is quite limited or repetitive.")
            actions.append("Experiment with more ambitious word choices.")

    # LINKING WORDS
    LINKERS = [
        "however", "therefore", "in addition", "furthermore", "moreover",
        "for example", "for instance", "consequently", "as a result"
    ]

    if any(l in text for l in LINKERS):
        strengths.append("You use linking words effectively to connect ideas.")
    else:
        weaknesses.append("Your writing lacks linking words to guide the reader.")
        actions.append("Use phrases like 'however', 'in addition', or 'for example' to show connections.")

    # ARGUMENT STRUCTURE
    ARG_MARKERS = ["because", "this shows", "this suggests", "therefore", "as a result"]

    if any(m in text for m in ARG_MARKERS):
        strengths.append("You explain your points with reasoning or evidence.")
    else:
        weaknesses.append("Some points are stated without explanation.")
        actions.append("After making a point, add a phrase like 'this shows that…' to explain it.")

    # NEXT STEPS
    next_steps.append("Read your work aloud to check that it flows logically.")
    next_steps.append("Compare your structure to a model answer to see how you could improve organisation.")
    next_steps.append("Ask a peer or teacher to highlight one unclear section, then rewrite it.")

    return strengths, weaknesses, actions, next_steps

# -----------------------------
# UI
# -----------------------------
uploaded = st.file_uploader(
    "Upload a document (.docx, .xlsx, .pptx)",
    type=["docx", "xlsx", "pptx"]
)

if st.button("🔄 Reset"):
    st.session_state.clear()
    st.rerun()

if uploaded:
    ext = os.path.splitext(uploaded.name)[1].lower()

    st.info(f"File detected: **{uploaded.name}** ({ext})")

    strengths, weaknesses, actions, next_steps = run_swan_analysis(uploaded, ext)

    st.subheader("Strengths")
    for s in strengths[:5]:
        st.write("•", s)

    st.subheader("Weaknesses")
    for w in weaknesses[:5]:
        st.write("•", w)

    st.subheader("Actions")
    for a in actions[:5]:
        st.write("•", a)

    st.subheader("Next Steps")
    for n in next_steps[:5]:
        st.write("•", n)

    st.divider()

    report_text = f"SWAN Feedback Report: {uploaded.name}\n"
    report_text += "="*30 + "\n\n"
    report_text += "STRENGTHS:\n" + "\n".join([f"- {s}" for s in strengths]) + "\n\n"
    report_text += "WEAKNESSES:\n" + "\n".join([f"- {w}" for w in weaknesses]) + "\n\n"
    report_text += "ACTIONS:\n" + "\n".join([f"- {a}" for a in actions]) + "\n\n"
    report_text += "NEXT STEPS:\n" + "\n".join([f"- {n}" for n in next_steps])

    st.download_button(
        label="📥 Download Feedback as Text",
        data=report_text,
        file_name=f"SWAN_Feedback_{uploaded.name}.txt",
        mime="text/plain"
    )

else:
    st.caption("Waiting for a file upload to begin SWAN analysis…")